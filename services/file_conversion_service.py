"""
File Conversion Service for the File Viewer & Editor System

This service handles conversion of various file types to web-compatible formats
for viewing and editing in the browser.
"""

import os
import mimetypes
from typing import Dict, Any, Optional, Union
from pathlib import Path
import logging

# Document processing imports
try:
    from docx import Document
    from docx.shared import Inches
except ImportError:
    Document = None

try:
    import openpyxl
    from openpyxl.utils import get_column_letter
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from PIL import Image, ExifTags
except ImportError:
    Image = None
    ExifTags = None

import bleach

logger = logging.getLogger(__name__)


class FileConversionService:
    """Service for converting files to web-compatible formats"""
    
    def __init__(self):
        self.supported_document_types = {
            '.docx': self._convert_docx_to_html,
            '.doc': self._convert_docx_to_html,  # Will need additional handling
            '.xlsx': self._convert_xlsx_to_html,
            '.xls': self._convert_xlsx_to_html,   # Will need additional handling
            '.pptx': self._convert_pptx_to_html,
            '.ppt': self._convert_pptx_to_html,   # Will need additional handling
        }
        
        self.supported_pdf_types = {
            '.pdf': self._extract_pdf_content
        }
        
        self.supported_image_types = {
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp'
        }
    
    def can_convert(self, file_path: str) -> bool:
        """Check if a file can be converted by this service"""
        ext = Path(file_path).suffix.lower()
        return (ext in self.supported_document_types or 
                ext in self.supported_pdf_types or
                ext in self.supported_image_types)
    
    def convert_document_to_html(self, file_path: str) -> Dict[str, Any]:
        """
        Convert a document file to HTML format
        
        Args:
            file_path: Path to the document file
            
        Returns:
            Dict containing HTML content and metadata
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        ext = Path(file_path).suffix.lower()
        
        if ext not in self.supported_document_types:
            raise ValueError(f"Unsupported document type: {ext}")
        
        try:
            converter_func = self.supported_document_types[ext]
            result = converter_func(file_path)
            
            # Sanitize HTML content
            if 'html' in result:
                result['html'] = self._sanitize_html(result['html'])
            
            return result
            
        except Exception as e:
            logger.error(f"Error converting document {file_path}: {str(e)}")
            raise RuntimeError(f"Document conversion failed: {str(e)}")
    
    def extract_pdf_content(self, file_path: str) -> Dict[str, Any]:
        """
        Extract content and metadata from PDF file
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Dict containing extracted content and metadata
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        if not PyPDF2:
            raise RuntimeError("PyPDF2 library not available for PDF processing")
        
        try:
            return self._extract_pdf_content(file_path)
        except Exception as e:
            logger.error(f"Error extracting PDF content {file_path}: {str(e)}")
            raise RuntimeError(f"PDF extraction failed: {str(e)}")
    
    def get_image_metadata(self, file_path: str) -> Dict[str, Any]:
        """
        Extract metadata from image file
        
        Args:
            file_path: Path to the image file
            
        Returns:
            Dict containing image metadata
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        ext = Path(file_path).suffix.lower()
        
        if ext not in self.supported_image_types:
            raise ValueError(f"Unsupported image type: {ext}")
        
        if not Image:
            raise RuntimeError("Pillow library not available for image processing")
        
        try:
            return self._extract_image_metadata(file_path)
        except Exception as e:
            logger.error(f"Error extracting image metadata {file_path}: {str(e)}")
            raise RuntimeError(f"Image metadata extraction failed: {str(e)}")
    
    def _convert_docx_to_html(self, file_path: str) -> Dict[str, Any]:
        """Convert DOCX file to HTML"""
        if not Document:
            raise RuntimeError("python-docx library not available")
        
        doc = Document(file_path)
        html_content = []
        
        # Convert paragraphs to HTML
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Basic formatting - can be enhanced
                html_content.append(f"<p>{paragraph.text}</p>")
        
        # Extract basic metadata
        core_props = doc.core_properties
        metadata = {
            'title': core_props.title or '',
            'author': core_props.author or '',
            'created': core_props.created.isoformat() if core_props.created else None,
            'modified': core_props.modified.isoformat() if core_props.modified else None,
            'pages': len(doc.paragraphs),  # Approximation
        }
        
        return {
            'html': '\n'.join(html_content),
            'metadata': metadata,
            'type': 'document'
        }
    
    def _convert_xlsx_to_html(self, file_path: str) -> Dict[str, Any]:
        """Convert XLSX file to HTML table format"""
        if not openpyxl:
            raise RuntimeError("openpyxl library not available")
        
        workbook = openpyxl.load_workbook(file_path, data_only=True)
        sheets_html = {}
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            
            # Convert sheet to HTML table
            html_rows = []
            for row in sheet.iter_rows(values_only=True):
                if any(cell is not None for cell in row):  # Skip empty rows
                    cells = [f"<td>{cell if cell is not None else ''}</td>" for cell in row]
                    html_rows.append(f"<tr>{''.join(cells)}</tr>")
            
            if html_rows:
                sheets_html[sheet_name] = f"<table class='excel-table'>{''.join(html_rows)}</table>"
        
        metadata = {
            'sheets': list(workbook.sheetnames),
            'sheet_count': len(workbook.sheetnames),
        }
        
        return {
            'html': sheets_html,
            'metadata': metadata,
            'type': 'spreadsheet'
        }
    
    def _convert_pptx_to_html(self, file_path: str) -> Dict[str, Any]:
        """Convert PPTX file to HTML slides"""
        if not Presentation:
            raise RuntimeError("python-pptx library not available")
        
        prs = Presentation(file_path)
        slides_html = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(f"<p>{shape.text}</p>")
            
            if slide_content:
                slides_html.append({
                    'slide_number': i + 1,
                    'html': '\n'.join(slide_content)
                })
        
        metadata = {
            'slide_count': len(prs.slides),
            'title': prs.core_properties.title or '',
            'author': prs.core_properties.author or '',
        }
        
        return {
            'html': slides_html,
            'metadata': metadata,
            'type': 'presentation'
        }
    
    def _extract_pdf_content(self, file_path: str) -> Dict[str, Any]:
        """Extract text content from PDF"""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            pages_content = []
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    text = page.extract_text()
                    pages_content.append({
                        'page_number': page_num + 1,
                        'text': text
                    })
                except Exception as e:
                    logger.warning(f"Could not extract text from page {page_num + 1}: {str(e)}")
                    pages_content.append({
                        'page_number': page_num + 1,
                        'text': ''
                    })
            
            # Extract metadata
            metadata = {
                'page_count': len(pdf_reader.pages),
                'title': '',
                'author': '',
                'creator': '',
                'producer': '',
            }
            
            if pdf_reader.metadata:
                metadata.update({
                    'title': pdf_reader.metadata.get('/Title', ''),
                    'author': pdf_reader.metadata.get('/Author', ''),
                    'creator': pdf_reader.metadata.get('/Creator', ''),
                    'producer': pdf_reader.metadata.get('/Producer', ''),
                })
            
            return {
                'pages': pages_content,
                'metadata': metadata,
                'type': 'pdf'
            }
    
    def _extract_image_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from image file"""
        with Image.open(file_path) as img:
            metadata = {
                'width': img.width,
                'height': img.height,
                'format': img.format,
                'mode': img.mode,
                'size_bytes': os.path.getsize(file_path),
            }
            
            # Extract EXIF data if available
            if hasattr(img, '_getexif') and img._getexif():
                exif_data = {}
                exif = img._getexif()
                
                for tag_id, value in exif.items():
                    tag = ExifTags.TAGS.get(tag_id, tag_id)
                    exif_data[tag] = value
                
                metadata['exif'] = exif_data
            
            return {
                'metadata': metadata,
                'type': 'image'
            }
    
    def _sanitize_html(self, html_content: str) -> str:
        """Sanitize HTML content to prevent XSS attacks"""
        allowed_tags = [
            'p', 'br', 'strong', 'em', 'u', 'h1', 'h2', 'h3', 'h4', 'h5', 'h6',
            'ul', 'ol', 'li', 'table', 'tr', 'td', 'th', 'thead', 'tbody',
            'div', 'span', 'img'
        ]
        
        allowed_attributes = {
            'img': ['src', 'alt', 'width', 'height'],
            'table': ['class'],
            'td': ['colspan', 'rowspan'],
            'th': ['colspan', 'rowspan'],
        }
        
        return bleach.clean(
            html_content,
            tags=allowed_tags,
            attributes=allowed_attributes,
            strip=True
        )